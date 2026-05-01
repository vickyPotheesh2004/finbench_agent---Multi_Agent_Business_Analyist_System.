"""
N01 PDF Ingestor - Multi-Format Document Ingestion
PDR-BAAAI-001 Rev 1.0 Node N01

Wires N01b ImageProcessor (opt-in) for OCR + chart vision extraction.

CHANGELOG:
  2026-04-30 S9  Bug #2: rewrote _ingest_html() to extract headings and
                 table cells properly. Was: soup.get_text() only, returning
                 0 headings + 0 table_cells. Now: parses <h1>-<h6>, <table>,
                 <th>/<td>, <b>/<strong>, with page-offset estimation.
                 This is the highest-impact fix in the campaign.
"""

from __future__ import annotations

import csv
import json
import logging
import os
import re
from typing import Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".xlsx", ".csv",
    ".pptx", ".html", ".htm", ".txt", ".eml",
    ".json", ".xml",
}

HEADING_FONT_SIZE_MIN = 13.0

# Approx chars per visual page in a SEC 10-K HTML — used for page estimation
HTML_CHARS_PER_PAGE = 3000

SEC_SECTIONS = [
    "business", "risk factors", "properties",
    "legal proceedings", "mine safety",
    "market", "selected financial data",
    "management", "financial statements",
    "quantitative", "controls", "other information",
    "directors", "executive compensation",
    "security ownership", "certain relationships",
    "principal accountant",
]


class TableCell:
    """Structured representation of a single table cell."""
    __slots__ = ("row_header", "col_header", "value",
                 "page", "table_number", "section")

    def __init__(
        self,
        row_header:   str = "",
        col_header:   str = "",
        value:        str = "",
        page:         int = 0,
        table_number: int = 0,
        section:      str = "",
    ) -> None:
        self.row_header   = row_header
        self.col_header   = col_header
        self.value        = value
        self.page         = page
        self.table_number = table_number
        self.section      = section

    def to_dict(self) -> Dict:
        return {
            "row_header":   self.row_header,
            "col_header":   self.col_header,
            "value":        self.value,
            "page":         self.page,
            "table_number": self.table_number,
            "section":      self.section,
        }


class PDFIngestor:
    """
    N01 Multi-Format Document Ingestor.

    Two usage modes:
        1. ingestor.ingest(file_path) -> dict
        2. ingestor.run(ba_state)     -> BAState

    Optional N01b image processing (opt-in via enable_images=True):
        - OCR on scanned pages via pytesseract
        - Chart/infographic value extraction via Gemma4 multimodal
    """

    def __init__(
        self,
        enable_images: bool = False,
        llm_client          = None,
    ) -> None:
        self.enable_images = enable_images
        self._llm          = llm_client

    def run(self, state) -> object:
        """LangGraph N01 node entry point."""
        doc_path = getattr(state, "document_path", "") or ""

        if not doc_path:
            logger.warning("N01: no document_path in state")
            return state

        if not os.path.exists(doc_path):
            logger.error("N01: file not found - %s", doc_path)
            return state

        result = self.ingest(doc_path)

        state.raw_text          = result.get("raw_text",          "")
        state.table_cells       = result.get("table_cells",       [])
        state.heading_positions = result.get("heading_positions", [])

        if not state.company_name:
            state.company_name = result.get("company_name", "")
        if not state.doc_type:
            state.doc_type     = result.get("doc_type",     "")
        if not state.fiscal_year:
            state.fiscal_year  = result.get("fiscal_year",  "")

        logger.info(
            "N01 Ingestor: %d chars | %d table_cells | %d headings | %s",
            len(state.raw_text),
            len(state.table_cells),
            len(state.heading_positions),
            os.path.basename(doc_path),
        )

        # N01b - Image Processor (opt-in)
        if self.enable_images and doc_path.lower().endswith(".pdf"):
            try:
                from src.ingestion.image_processor import ImageProcessor
                image_proc = ImageProcessor(
                    enable_ocr    = True,
                    enable_vision = self._llm is not None,
                    llm_client    = self._llm,
                )
                state = image_proc.run(state)
                logger.info("N01b: image processing complete")
            except Exception as exc:
                logger.warning("N01b image processor failed: %s", exc)

        return state

    def ingest(self, file_path: str) -> Dict:
        """Ingest any supported document type."""
        ext = os.path.splitext(file_path)[1].lower()

        if ext == ".pdf":
            return self._ingest_pdf(file_path)
        elif ext in (".docx", ".doc"):
            return self._ingest_docx(file_path)
        elif ext in (".xlsx", ".xls"):
            return self._ingest_xlsx(file_path)
        elif ext == ".csv":
            return self._ingest_csv(file_path)
        elif ext == ".pptx":
            return self._ingest_pptx(file_path)
        elif ext in (".html", ".htm"):
            return self._ingest_html(file_path)
        elif ext == ".json":
            return self._ingest_json(file_path)
        else:
            return self._ingest_txt(file_path)

    def _ingest_pdf(self, file_path: str) -> Dict:
        raw_text          = ""
        table_cells       = []
        heading_positions = []

        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, start=1):
                    page_text = page.extract_text() or ""
                    raw_text += page_text + "\n"

                    tables = page.extract_tables() or []
                    for table_num, table in enumerate(tables, start=1):
                        if not table or len(table) < 2:
                            continue
                        col_headers = [
                            str(c).strip() if c else ""
                            for c in table[0]
                        ]
                        for row in table[1:]:
                            if not row:
                                continue
                            row_header = str(row[0]).strip() if row[0] else ""
                            for col_idx, cell_val in enumerate(row[1:], start=1):
                                col_header = (
                                    col_headers[col_idx]
                                    if col_idx < len(col_headers) else ""
                                )
                                value = str(cell_val).strip() if cell_val else ""
                                if value:
                                    table_cells.append(TableCell(
                                        row_header   = row_header,
                                        col_header   = col_header,
                                        value        = value,
                                        page         = page_num,
                                        table_number = table_num,
                                    ).to_dict())
        except ImportError:
            logger.warning("pdfplumber not installed - PDF tables skipped")
        except Exception as exc:
            logger.warning("pdfplumber error on %s: %s", file_path, exc)

        try:
            import fitz
            doc = fitz.open(file_path)
            for page_num, page in enumerate(doc, start=1):
                blocks = page.get_text("dict").get("blocks", [])
                for block in blocks:
                    for line in block.get("lines", []):
                        for span in line.get("spans", []):
                            font_size = span.get("size", 0)
                            flags     = span.get("flags", 0)
                            text      = span.get("text", "").strip()
                            is_bold   = bool(flags & 2 ** 4)
                            if font_size >= HEADING_FONT_SIZE_MIN and text and len(text) > 3:
                                heading_positions.append({
                                    "text":      text,
                                    "font_size": round(font_size, 1),
                                    "is_bold":   is_bold,
                                    "page":      page_num,
                                })
            doc.close()
        except ImportError:
            logger.warning("PyMuPDF not installed - headings skipped")
        except Exception as exc:
            logger.warning("PyMuPDF error on %s: %s", file_path, exc)

        company_name, doc_type, fiscal_year = self._extract_metadata(
            raw_text, heading_positions
        )
        return {
            "raw_text":          raw_text,
            "table_cells":       table_cells,
            "heading_positions": heading_positions,
            "company_name":      company_name,
            "doc_type":          doc_type,
            "fiscal_year":       fiscal_year,
        }

    def _ingest_docx(self, file_path: str) -> Dict:
        raw_text          = ""
        table_cells       = []
        heading_positions = []

        try:
            from docx import Document
            doc = Document(file_path)

            for para in doc.paragraphs:
                raw_text += para.text + "\n"
                if para.style and "Heading" in para.style.name:
                    level = para.style.name.replace("Heading ", "").strip()
                    heading_positions.append({
                        "text":      para.text.strip(),
                        "font_size": 16.0 if level == "1" else 13.0,
                        "is_bold":   True,
                        "page":      0,
                    })

            for t_num, table in enumerate(doc.tables, start=1):
                rows = list(table.rows)
                if len(rows) < 2:
                    continue
                col_headers = [c.text.strip() for c in rows[0].cells]
                for row in rows[1:]:
                    cells      = row.cells
                    row_header = cells[0].text.strip() if cells else ""
                    for col_idx, cell in enumerate(cells[1:], start=1):
                        col_header = (
                            col_headers[col_idx]
                            if col_idx < len(col_headers) else ""
                        )
                        value = cell.text.strip()
                        if value:
                            table_cells.append(TableCell(
                                row_header   = row_header,
                                col_header   = col_header,
                                value        = value,
                                page         = 0,
                                table_number = t_num,
                            ).to_dict())
        except ImportError:
            logger.warning("python-docx not installed")
        except Exception as exc:
            logger.warning("DOCX error: %s", exc)

        company_name, doc_type, fiscal_year = self._extract_metadata(
            raw_text, heading_positions
        )
        return {
            "raw_text": raw_text, "table_cells": table_cells,
            "heading_positions": heading_positions,
            "company_name": company_name, "doc_type": doc_type,
            "fiscal_year": fiscal_year,
        }

    def _ingest_xlsx(self, file_path: str) -> Dict:
        raw_text    = ""
        table_cells = []

        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            for sheet_name in wb.sheetnames:
                ws   = wb[sheet_name]
                rows = list(ws.iter_rows(values_only=True))
                if not rows:
                    continue
                col_headers = [
                    str(c).strip() if c is not None else ""
                    for c in rows[0]
                ]
                raw_text += f"\n[Sheet: {sheet_name}]\n"
                for row in rows[1:]:
                    row_header = str(row[0]).strip() if row[0] is not None else ""
                    for col_idx, val in enumerate(row[1:], start=1):
                        if val is not None:
                            col_header = (
                                col_headers[col_idx]
                                if col_idx < len(col_headers) else ""
                            )
                            value = str(val).strip()
                            raw_text += f"{row_header} | {col_header} | {value}\n"
                            table_cells.append(TableCell(
                                row_header   = row_header,
                                col_header   = col_header,
                                value        = value,
                                page         = 0,
                                table_number = 0,
                                section      = sheet_name,
                            ).to_dict())
        except ImportError:
            logger.warning("openpyxl not installed")
        except Exception as exc:
            logger.warning("XLSX error: %s", exc)

        return {
            "raw_text": raw_text, "table_cells": table_cells,
            "heading_positions": [],
            "company_name": "", "doc_type": "", "fiscal_year": "",
        }

    def _ingest_csv(self, file_path: str) -> Dict:
        raw_text    = ""
        table_cells = []

        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                reader = csv.reader(f)
                rows   = list(reader)
            if rows:
                col_headers = rows[0]
                for row_idx, row in enumerate(rows[1:], start=1):
                    row_header = row[0].strip() if row else ""
                    for col_idx, val in enumerate(row[1:], start=1):
                        col_header = (
                            col_headers[col_idx].strip()
                            if col_idx < len(col_headers) else ""
                        )
                        value = val.strip()
                        raw_text += f"{row_header} | {col_header} | {value}\n"
                        if value:
                            table_cells.append(TableCell(
                                row_header   = row_header,
                                col_header   = col_header,
                                value        = value,
                                page         = 0,
                                table_number = row_idx,
                            ).to_dict())
        except Exception as exc:
            logger.warning("CSV error: %s", exc)

        return {
            "raw_text": raw_text, "table_cells": table_cells,
            "heading_positions": [],
            "company_name": "", "doc_type": "", "fiscal_year": "",
        }

    def _ingest_pptx(self, file_path: str) -> Dict:
        raw_text          = ""
        heading_positions = []

        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides, start=1):
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                raw_text += text + "\n"
                                if para.runs and para.runs[0].font.size:
                                    font_pt = para.runs[0].font.size.pt
                                    if font_pt >= HEADING_FONT_SIZE_MIN:
                                        heading_positions.append({
                                            "text":      text,
                                            "font_size": round(font_pt, 1),
                                            "is_bold":   bool(para.runs[0].font.bold),
                                            "page":      slide_num,
                                        })
        except ImportError:
            logger.warning("python-pptx not installed")
        except Exception as exc:
            logger.warning("PPTX error: %s", exc)

        return {
            "raw_text": raw_text, "table_cells": [],
            "heading_positions": heading_positions,
            "company_name": "", "doc_type": "", "fiscal_year": "",
        }

    # ════════════════════════════════════════════════════════════════════════
    # HTML INGESTION  ── BUG #2 FIX LIVES HERE ──
    # ════════════════════════════════════════════════════════════════════════

    def _ingest_html(self, file_path: str) -> Dict:
        """Bug #2 fix: full HTML extraction.

        Old behaviour: soup.get_text() only → 0 headings, 0 tables.
        New behaviour: extracts <h1>-<h6>, <table>/<th>/<td>, <b>/<strong>,
                       and estimates page numbers from byte offset.
        """
        raw_text          = ""
        table_cells       = []
        heading_positions = []

        try:
            # Suppress XML-as-HTML warning (SEC files are valid for our purpose)
            import warnings
            try:
                from bs4 import XMLParsedAsHTMLWarning
                warnings.filterwarnings("ignore", category=XMLParsedAsHTMLWarning)
            except ImportError:
                pass

            from bs4 import BeautifulSoup

            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()

            # Try lxml first, fall back to html.parser
            try:
                soup = BeautifulSoup(content, "lxml")
            except Exception:
                soup = BeautifulSoup(content, "html.parser")

            # ── Step 1: extract clean text ────────────────────────────────
            # Remove script and style first so they don't leak into raw_text
            for tag in soup(["script", "style", "noscript", "meta", "link"]):
                tag.decompose()

            raw_text = soup.get_text(separator="\n")
            # Collapse 3+ newlines to 2 (preserve paragraph breaks)
            raw_text = re.sub(r"\n{3,}", "\n\n", raw_text)
            # Collapse multiple spaces but keep newlines
            raw_text = re.sub(r"[ \t]+", " ", raw_text)

            # ── Step 2: extract headings <h1>-<h6> ────────────────────────
            heading_positions.extend(
                self._html_extract_headings(soup, raw_text)
            )

            # ── Step 3: extract bold-as-heading <b>, <strong> ────────────
            heading_positions.extend(
                self._html_extract_bold_headings(soup, raw_text)
            )

            # ── Step 4: extract table cells ───────────────────────────────
            table_cells.extend(
                self._html_extract_table_cells(soup, raw_text)
            )

            # ── Step 5: dedupe headings (some 10-Ks repeat) ───────────────
            heading_positions = self._dedupe_headings(heading_positions)

            logger.debug(
                "HTML extracted: %d chars, %d headings, %d table cells",
                len(raw_text), len(heading_positions), len(table_cells),
            )

        except ImportError:
            # bs4 missing — fall back to crude regex strip
            logger.warning("BeautifulSoup not installed — using regex fallback")
            try:
                with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                    raw_text = f.read()
                raw_text = re.sub(r"<script[^>]*>.*?</script>", " ",
                                  raw_text, flags=re.DOTALL | re.IGNORECASE)
                raw_text = re.sub(r"<style[^>]*>.*?</style>", " ",
                                  raw_text, flags=re.DOTALL | re.IGNORECASE)
                raw_text = re.sub(r"<[^>]+>", " ", raw_text)
                raw_text = re.sub(r"\s+", " ", raw_text)
            except Exception as exc:
                logger.warning("HTML regex fallback error: %s", exc)
        except Exception as exc:
            logger.warning("HTML error: %s", exc)

        company_name, doc_type, fiscal_year = self._extract_metadata(
            raw_text, heading_positions
        )

        return {
            "raw_text":          raw_text,
            "table_cells":       table_cells,
            "heading_positions": heading_positions,
            "company_name":      company_name,
            "doc_type":          doc_type,
            "fiscal_year":       fiscal_year,
        }

    @staticmethod
    def _estimate_page_from_offset(offset: int) -> int:
        """Estimate visual page number from byte offset in HTML text."""
        if offset <= 0:
            return 1
        return max(1, (offset // HTML_CHARS_PER_PAGE) + 1)

    def _html_extract_headings(self, soup, raw_text: str) -> List[Dict]:
        """Extract <h1>-<h6> elements as heading_positions."""
        out = []
        # font_size by level: h1=20, h2=18, h3=16, h4=14, h5=13, h6=13
        size_by_level = {1: 20.0, 2: 18.0, 3: 16.0, 4: 14.0, 5: 13.0, 6: 13.0}

        for level in range(1, 7):
            for tag in soup.find_all(f"h{level}"):
                text = tag.get_text(separator=" ", strip=True)
                if not text or len(text) < 3 or len(text) > 200:
                    continue

                # Estimate page by finding text in raw_text
                page = 1
                snippet = text[:30]
                idx = raw_text.find(snippet)
                if idx >= 0:
                    page = self._estimate_page_from_offset(idx)

                out.append({
                    "text":      text,
                    "font_size": size_by_level[level],
                    "is_bold":   True,
                    "page":      page,
                })
        return out

    def _html_extract_bold_headings(self, soup, raw_text: str) -> List[Dict]:
        """Extract <b>/<strong> short text as candidate headings.

        SEC 10-Ks often use bold text as section markers without <hN> tags.
        We treat <b>/<strong> with 3-100 chars and 1-12 words as heading-like.
        """
        out = []
        for tag in soup.find_all(["b", "strong"]):
            text = tag.get_text(separator=" ", strip=True)
            if not text:
                continue
            wc = len(text.split())
            # Reject sentences (likely mid-paragraph emphasis)
            if len(text) < 3 or len(text) > 120 or wc > 12:
                continue
            # Reject pure numbers / page refs
            if re.match(r"^[\d\s,.\-$()]+$", text):
                continue

            page = 1
            snippet = text[:30]
            idx = raw_text.find(snippet)
            if idx >= 0:
                page = self._estimate_page_from_offset(idx)

            out.append({
                "text":      text,
                "font_size": 13.5,    # H2-borderline
                "is_bold":   True,
                "page":      page,
            })
        return out

    def _html_extract_table_cells(self, soup, raw_text: str) -> List[Dict]:
        """Extract every <td>/<th> cell with row+column headers from <table>."""
        cells: List[Dict] = []

        for t_num, table in enumerate(soup.find_all("table"), start=1):
            rows = table.find_all("tr")
            if len(rows) < 2:
                continue

            # Estimate page from table position
            page = 1
            first_row_text = (
                rows[0].get_text(" ", strip=True) if rows else ""
            )
            if first_row_text:
                idx = raw_text.find(first_row_text[:50])
                if idx >= 0:
                    page = self._estimate_page_from_offset(idx)

            # Detect header row (uses <th>, OR first row with all bold cells,
            # OR just take row 0 as header)
            header_row = None
            for r in rows:
                ths = r.find_all("th")
                if ths:
                    header_row = r
                    break
            if header_row is None:
                header_row = rows[0]

            col_headers = [
                c.get_text(" ", strip=True)
                for c in header_row.find_all(["th", "td"])
            ]

            # Determine data rows
            try:
                start_idx = rows.index(header_row) + 1
            except ValueError:
                start_idx = 1
            data_rows = rows[start_idx:]

            for row in data_rows:
                row_cells = row.find_all(["th", "td"])
                if not row_cells:
                    continue

                row_header = row_cells[0].get_text(" ", strip=True)
                for col_idx, c in enumerate(row_cells[1:], start=1):
                    value = c.get_text(" ", strip=True)
                    # Skip empty cells, skip pure punctuation noise
                    if not value or value in ("—", "-", "$", "(", ")"):
                        continue
                    col_header = (
                        col_headers[col_idx]
                        if col_idx < len(col_headers) else ""
                    )
                    cells.append(TableCell(
                        row_header   = row_header,
                        col_header   = col_header,
                        value        = value,
                        page         = page,
                        table_number = t_num,
                    ).to_dict())

        return cells

    @staticmethod
    def _dedupe_headings(headings: List[Dict]) -> List[Dict]:
        """Deduplicate by text (case-insensitive). Keep first occurrence."""
        seen = set()
        out  = []
        for h in headings:
            key = (h.get("text", "") or "").strip().lower()
            if not key or key in seen:
                continue
            seen.add(key)
            out.append(h)
        return out

    def _ingest_txt(self, file_path: str) -> Dict:
        raw_text = ""
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                raw_text = f.read()
        except Exception as exc:
            logger.warning("TXT error: %s", exc)

        return {
            "raw_text": raw_text, "table_cells": [],
            "heading_positions": [],
            "company_name": "", "doc_type": "", "fiscal_year": "",
        }

    def _ingest_json(self, file_path: str) -> Dict:
        raw_text = ""
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                data     = json.load(f)
            raw_text = json.dumps(data, indent=2)
        except Exception as exc:
            logger.warning("JSON error: %s", exc)

        return {
            "raw_text": raw_text, "table_cells": [],
            "heading_positions": [],
            "company_name": "", "doc_type": "", "fiscal_year": "",
        }

    def _extract_metadata(
        self,
        raw_text:          str,
        heading_positions: List[Dict],
    ) -> Tuple[str, str, str]:
        company_name = self._extract_company(raw_text)
        doc_type     = self._extract_doc_type(raw_text)
        fiscal_year  = self._extract_fiscal_year(raw_text)
        return company_name, doc_type, fiscal_year

    @staticmethod
    def _extract_company(text: str) -> str:
        snippet  = text[:2000]
        patterns = [
            r"((?:[A-Z][a-z]+\s*){1,4}(?:Inc|Corp|Ltd|LLC|Co|Company|Group|Holdings)\.?)",
        ]
        for pattern in patterns:
            m = re.search(pattern, snippet)
            if m:
                return m.group(1).strip()
        return ""

    @staticmethod
    def _extract_doc_type(text: str) -> str:
        snippet = text[:3000].upper()
        for dtype in ["10-K", "10-Q", "8-K", "S-1", "20-F", "6-K", "DEF 14A"]:
            if dtype.replace("-", "") in snippet.replace("-", "").replace(" ", ""):
                return dtype
        return "UNKNOWN"

    @staticmethod
    def _extract_fiscal_year(text: str) -> str:
        patterns = [
            r"[Ff]iscal [Yy]ear\s*(?:ended|ending)?\s*(?:(?:September|December|March|June)\s+\d{1,2},?\s*)?(\d{4})",
            r"[Ff][Yy]\s*(\d{4})",
            r"[Yy]ear [Ee]nded\s+\w+\s+\d{1,2},?\s*(\d{4})",
            r"[Aa]nnual [Rr]eport\s+(\d{4})",
        ]
        for pattern in patterns:
            m = re.search(pattern, text[:5000])
            if m:
                return f"FY{m.group(1)}"
        return ""


def run_pdf_ingestor(state, enable_images: bool = False, llm_client=None) -> object:
    """Convenience wrapper for LangGraph N01 node."""
    return PDFIngestor(
        enable_images = enable_images,
        llm_client    = llm_client,
    ).run(state)